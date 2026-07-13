"""Structure-aware attachment chunking (issue #89).

Attachment text used to be handed to the same generic ``SentenceSplitter`` as email
bodies. That splitter needs *sentence boundaries* to find split points — but a
flattened spreadsheet/table has none, so it emitted ONE giant chunk and bge-m3
truncated it at 8192 tokens, silently losing the tail rows (observed in the v2
rebuild: ``Token indices sequence length is longer than the specified maximum
(30830 > 8192)``, see #87).

The fix, implemented here, is a **per-MIME chunker** that splits each attachment
using a splitter that understands its FORMAT, *before* embedding, so every emitted
chunk is a coherent, embedder-sized unit:

* **Spreadsheets (.xlsx/.csv)** — parsed with **pandas**; chunked by **row-groups**,
  **repeating the header row in every chunk** so each chunk is a self-describing
  mini-table. One chunk-stream per sheet/tab. Rows are never cut mid-row.
* **PDF** — split by **page** (a page that alone exceeds the budget is prose-split).
* **PPTX** — split by **slide** (an over-budget slide is prose-split).
* **DOCX** — split by **heading/section** (falling back to paragraph groups).
* **Fallback** — the token-aware prose splitter for genuinely prose-like attachments
  or whenever structure parsing fails/library is missing.

Every emitted chunk is guaranteed ``token_len(chunk) <= budget`` (``budget`` targets
the profile ``chunk_size`` and is a HARD cap well under 8192). Because each chunk is
already a whole Document sized under the budget, the downstream ``SentenceSplitter``
in :func:`src.indexing.contextual_index.build_contextual_index` leaves it intact
(it never re-splits a document that is already under ``chunk_size``, and dedup keeps
distinct chunks distinct).

The parsing operates on the attachment's **raw bytes** (re-parsed here), not on the
already-flattened extracted text, because the extract handlers drop all structural
markers (sheet/page/slide boundaries, the header row). The flattened text is kept as
the fallback input for prose-like content and when a format library is unavailable.
"""

from __future__ import annotations

import io
import re
from typing import Callable, List, Optional

# Default token budget for an attachment chunk. Targets the typical profile
# ``chunk_size`` (512) but stays a hard cap regardless of the caller. Always well
# under the bge-m3 8192 ceiling, so no chunk can ever be truncated at embed time.
DEFAULT_CHUNK_BUDGET = 512


def _default_token_len(text: str) -> int:
    """A conservative, tokenizer-free token estimate.

    We deliberately *over*-count relative to a real sub-word tokenizer so that a
    chunk sized to ``budget`` here is comfortably under ``budget`` for bge-m3 — the
    whole point is to never trip the 8192 truncation, and over-counting only makes
    chunks a little smaller, never larger. Roughly one token per whitespace- or
    punctuation-delimited unit. Injectable via ``token_len`` for exact sizing/tests.
    """
    if not text:
        return 0
    # Split on whitespace and on runs of punctuation/symbols so that "210,000,000"
    # or "a\tb\tc" counts several units, mirroring how a sub-word tokenizer fragments
    # dense numeric/tabular text.
    return len(re.findall(r"\w+|[^\w\s]", text))


def _split_prose(text: str, budget: int, token_len: Callable[[str], int]) -> List[str]:
    """Fallback splitter for prose-like text: greedily pack lines into chunks whose
    token count stays under ``budget``. A single line longer than the budget is hard-
    wrapped on whitespace so no chunk ever exceeds the cap."""
    chunks: List[str] = []
    cur: List[str] = []
    cur_tokens = 0
    for line in text.split("\n"):
        lt = token_len(line)
        if lt > budget:
            # Flush the accumulator, then hard-wrap this over-long line on words.
            if cur:
                chunks.append("\n".join(cur))
                cur, cur_tokens = [], 0
            chunks.extend(_hard_wrap(line, budget, token_len))
            continue
        if cur_tokens + lt > budget and cur:
            chunks.append("\n".join(cur))
            cur, cur_tokens = [], 0
        cur.append(line)
        cur_tokens += lt
    if cur:
        chunks.append("\n".join(cur))
    return [c for c in chunks if c.strip()]


def _hard_wrap(line: str, budget: int, token_len: Callable[[str], int]) -> List[str]:
    """Wrap a single over-budget line on word boundaries so each piece fits."""
    words = line.split(" ")
    out: List[str] = []
    cur: List[str] = []
    for w in words:
        trial = " ".join([*cur, w]) if cur else w
        if token_len(trial) > budget and cur:
            out.append(" ".join(cur))
            cur = [w]
        else:
            cur.append(w)
    if cur:
        out.append(" ".join(cur))
    return out


def _row_str(values) -> str:
    """Render one spreadsheet row as a tab-joined line (matches the xlsx handler)."""
    return "\t".join("" if v is None else str(v) for v in values)


def _chunk_rows(
    header: str,
    rows: List[str],
    budget: int,
    token_len: Callable[[str], int],
    prefix: str = "",
) -> List[str]:
    """Pack data ``rows`` into chunks, repeating ``header`` (and optional ``prefix``,
    e.g. a sheet label) at the top of every chunk so each chunk is a self-describing
    mini-table. Never splits a row across chunks. A single row wider than the budget
    is emitted on its own (with the header) — it is still one coherent record and
    the budget is advisory for such pathological rows, but this keeps rows intact."""
    lead = [p for p in (prefix, header) if p]
    lead_tokens = sum(token_len(x) for x in lead)
    chunks: List[str] = []
    cur: List[str] = []
    cur_tokens = lead_tokens
    for r in rows:
        rt = token_len(r)
        if cur and cur_tokens + rt > budget:
            chunks.append("\n".join([*lead, *cur]))
            cur, cur_tokens = [], lead_tokens
        cur.append(r)
        cur_tokens += rt
    if cur:
        chunks.append("\n".join([*lead, *cur]))
    return chunks


def _chunk_spreadsheet(
    data: bytes, filename: str, mime: str, budget: int, token_len: Callable[[str], int]
) -> Optional[List[str]]:
    """Chunk a spreadsheet by row-groups, one stream per sheet, header repeated.

    Returns ``None`` (signalling "fall back") if pandas is missing or the bytes do
    not parse as a spreadsheet."""
    try:
        import pandas as pd
    except Exception:
        return None

    name = (filename or "").lower()
    is_csv = name.endswith(".csv") or "csv" in (mime or "").lower()
    try:
        if is_csv:
            # header=None: keep the first row as data too, so a headerless ledger's
            # first line is not silently promoted/lost. The first row still repeats
            # as the "header" line of every chunk (self-describing), which is the
            # right behaviour whether or not it is a real header.
            frames = {
                "": pd.read_csv(io.BytesIO(data), header=None, dtype=str, keep_default_na=False)
            }
        else:
            frames = pd.read_excel(
                io.BytesIO(data), sheet_name=None, header=None, dtype=str, engine="openpyxl"
            )
            frames = {k: v.fillna("") for k, v in frames.items()}
    except Exception:
        return None

    chunks: List[str] = []
    for sheet_name, df in frames.items():
        if df is None or df.empty:
            continue
        rows = [_row_str(row) for row in df.itertuples(index=False, name=None)]
        if not rows:
            continue
        header = rows[0]
        body = rows[1:]
        prefix = f"[sheet: {sheet_name}]" if sheet_name else ""
        if not body:
            # Header-only (single-row) sheet: one chunk carrying just that row.
            lead = [p for p in (prefix, header) if p]
            chunks.append("\n".join(lead))
            continue
        chunks.extend(_chunk_rows(header, body, budget, token_len, prefix=prefix))
    return chunks or None


def _chunk_pdf(data: bytes, budget: int, token_len: Callable[[str], int]) -> Optional[List[str]]:
    """Chunk a PDF by page; a page over the budget is prose-split into page-groups."""
    try:
        import pypdf
    except Exception:
        return None
    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = [(p.extract_text() or "") for p in reader.pages]
    except Exception:
        return None
    if not any(p.strip() for p in pages):
        return None
    chunks: List[str] = []
    for page in pages:
        page = page.strip()
        if not page:
            continue
        if token_len(page) > budget:
            chunks.extend(_split_prose(page, budget, token_len))
        else:
            chunks.append(page)
    return chunks or None


def _chunk_pptx(data: bytes, budget: int, token_len: Callable[[str], int]) -> Optional[List[str]]:
    """Chunk a PPTX by slide; an over-budget slide is prose-split."""
    try:
        import pptx
    except Exception:
        return None
    try:
        pres = pptx.Presentation(io.BytesIO(data))
    except Exception:
        return None
    chunks: List[str] = []
    for slide in pres.slides:
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        ]
        text = "\n".join(p for p in parts if p).strip()
        if not text:
            continue
        if token_len(text) > budget:
            chunks.extend(_split_prose(text, budget, token_len))
        else:
            chunks.append(text)
    return chunks or None


def _chunk_docx(data: bytes, budget: int, token_len: Callable[[str], int]) -> Optional[List[str]]:
    """Chunk a DOCX by heading/section, falling back to paragraph groups.

    A "section" is the run of paragraphs from one heading up to (but excluding) the
    next heading. Sections over the budget are prose-split; documents with no
    headings degrade to a single prose-split of the whole body."""
    try:
        import docx
    except Exception:
        return None
    try:
        d = docx.Document(io.BytesIO(data))
    except Exception:
        return None

    sections: List[List[str]] = []
    cur: List[str] = []
    for p in d.paragraphs:
        style = (getattr(p.style, "name", "") or "").lower()
        is_heading = style.startswith("heading") or style == "title"
        if is_heading and cur:
            sections.append(cur)
            cur = []
        if p.text.strip() or is_heading:
            cur.append(p.text)
    if cur:
        sections.append(cur)
    if not sections:
        return None

    chunks: List[str] = []
    for sec in sections:
        text = "\n".join(sec).strip()
        if not text:
            continue
        if token_len(text) > budget:
            chunks.extend(_split_prose(text, budget, token_len))
        else:
            chunks.append(text)
    return chunks or None


def chunk_attachment(
    *,
    data: bytes,
    text: str,
    mime: str,
    filename: str,
    budget: int = DEFAULT_CHUNK_BUDGET,
    token_len: Optional[Callable[[str], int]] = None,
) -> List[str]:
    """Split one attachment into embedder-sized, structure-aware chunks.

    Parameters
    ----------
    data:
        The attachment's raw bytes — re-parsed here so the format's own structural
        units (sheet rows, PDF pages, slides, DOCX sections) can drive the split.
    text:
        The already-extracted (flattened) text, used as the fallback input for prose
        splitting and when a format-aware parse is unavailable/fails.
    mime, filename:
        Used to route to the right per-format chunker.
    budget:
        Hard token cap per chunk (defaults to the profile ``chunk_size``). Every
        returned chunk satisfies ``token_len(chunk) <= budget`` except a single
        spreadsheet row wider than the budget, which is kept intact (never cut).
    token_len:
        Injectable token counter. Defaults to a conservative tokenizer-free estimate
        that over-counts vs. bge-m3, so a chunk sized here stays under the real
        8192-token embed ceiling.

    Returns
    -------
    list[str]
        One string per chunk, in reading order. Never empty when ``text`` is
        non-empty (falls back to a single/prose chunk if structure parsing yields
        nothing).
    """
    tl = token_len or _default_token_len
    name = (filename or "").lower()
    m = (mime or "").lower()

    structured: Optional[List[str]] = None
    if name.endswith((".xlsx", ".csv")) or "spreadsheetml" in m or "csv" in m:
        structured = _chunk_spreadsheet(data, filename, mime, budget, tl)
    elif name.endswith(".pdf") or "pdf" in m:
        structured = _chunk_pdf(data, budget, tl)
    elif name.endswith(".pptx") or "presentationml" in m:
        structured = _chunk_pptx(data, budget, tl)
    elif name.endswith(".docx") or "wordprocessingml" in m:
        structured = _chunk_docx(data, budget, tl)

    if structured:
        return structured

    # Fallback: prose-split the flattened extracted text. Guarantees non-empty
    # output for non-empty text and keeps every chunk under the budget.
    prose = _split_prose(text, budget, tl)
    return prose or ([text] if text.strip() else [])
