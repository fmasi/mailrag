"""PPTX handler (python-pptx)."""

from __future__ import annotations

import io

from src.attachments.extract.result import ExtractResult, Status, ok


class PptxHandler:
    def can_handle(self, mime: str, filename: str) -> bool:
        return (filename or "").lower().endswith(".pptx") or "presentationml" in (
            mime or ""
        ).lower()

    def extract(self, data: bytes, mime: str, filename: str) -> ExtractResult:
        try:
            import pptx
        except Exception:
            return ExtractResult("", Status.BINARY, "pptx")
        try:
            pres = pptx.Presentation(io.BytesIO(data))
            out = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return ok("\n".join(out), "pptx")
        except Exception:
            return ExtractResult("", Status.ERROR, "pptx")
