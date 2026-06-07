"""Best-effort, guarded text extraction from attachment bytes.

Every extractor catches its own import/parse errors and returns an ExtractResult
rather than raising, so a missing library or system tool degrades the status (the
caller still gets the raw file). See the 1a design spec.
"""
from __future__ import annotations

import io
from dataclasses import dataclass
from html.parser import HTMLParser


@dataclass
class ExtractResult:
    text: str
    status: str       # extracted | empty | binary | ocr_unavailable | error
    extractor: str


def _ok(text, extractor):
    return ExtractResult(text=text, status=("extracted" if text.strip() else "empty"),
                         extractor=extractor)


class _Stripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data):
        self.parts.append(data)


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _html(data, filename):
    p = _Stripper()
    p.feed(_decode(data))
    return _ok(" ".join(" ".join(p.parts).split()), "html.parser")


def _pdf(data, filename):
    try:
        import pypdf
    except Exception:
        return _ocr_pdf(data, filename)
    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        text = "\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception:
        return ExtractResult("", "error", "pypdf")
    if text.strip():
        return _ok(text, "pypdf")
    return _ocr_pdf(data, filename)   # no text layer -> try OCR


def _ocr_pdf(data, filename):
    try:
        import pdf2image
        import pytesseract
    except Exception:
        return ExtractResult("", "ocr_unavailable", "ocr")
    try:
        images = pdf2image.convert_from_bytes(data)
        text = "\n".join(pytesseract.image_to_string(im) for im in images)
        return _ok(text, "ocr")
    except Exception:
        return ExtractResult("", "ocr_unavailable", "ocr")


def _image_ocr(data, filename):
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return ExtractResult("", "ocr_unavailable", "ocr")
    try:
        text = pytesseract.image_to_string(Image.open(io.BytesIO(data)))
        return _ok(text, "ocr")
    except Exception:
        return ExtractResult("", "ocr_unavailable", "ocr")


def _docx(data, filename):
    try:
        import docx
    except Exception:
        return ExtractResult("", "binary", "python-docx")
    try:
        d = docx.Document(io.BytesIO(data))
        return _ok("\n".join(p.text for p in d.paragraphs), "python-docx")
    except Exception:
        return ExtractResult("", "error", "python-docx")


def _xlsx(data, filename):
    try:
        import openpyxl
    except Exception:
        return ExtractResult("", "binary", "openpyxl")
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                out.append("\t".join("" if c is None else str(c) for c in row))
        return _ok("\n".join(out), "openpyxl")
    except Exception:
        return ExtractResult("", "error", "openpyxl")


def _pptx(data, filename):
    try:
        import pptx
    except Exception:
        return ExtractResult("", "binary", "python-pptx")
    try:
        pres = pptx.Presentation(io.BytesIO(data))
        out = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return _ok("\n".join(out), "python-pptx")
    except Exception:
        return ExtractResult("", "error", "python-pptx")


def extract_text(data: bytes, mime: str, filename: str) -> ExtractResult:
    """Return extracted text + status for the given bytes; never raises."""
    m = (mime or "").lower()
    name = (filename or "").lower()
    try:
        if m in ("text/plain", "text/csv") or name.endswith((".txt", ".csv")):
            return _ok(_decode(data), "stdlib")
        if m == "text/calendar" or name.endswith(".ics"):
            return _ok(_decode(data), "stdlib")
        if m == "text/html" or name.endswith((".html", ".htm")):
            return _html(data, filename)
        if m == "application/pdf" or name.endswith(".pdf"):
            return _pdf(data, filename)
        if name.endswith(".docx") or "wordprocessingml" in m:
            return _docx(data, filename)
        if name.endswith(".xlsx") or "spreadsheetml" in m:
            return _xlsx(data, filename)
        if name.endswith(".pptx") or "presentationml" in m:
            return _pptx(data, filename)
        if m.startswith("image/"):
            return _image_ocr(data, filename)
    except Exception:
        return ExtractResult("", "error", "extract")
    return ExtractResult("", "binary", "none")
