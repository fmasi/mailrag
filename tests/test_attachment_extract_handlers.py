import io
import unittest
from unittest import mock
from src.attachments.extract.handlers.plaintext import PlaintextHandler
from src.attachments.extract.handlers.html import HtmlHandler
from src.attachments.extract.handlers.docx import DocxHandler
from src.attachments.extract.handlers.xlsx import XlsxHandler
from src.attachments.extract.handlers.pptx import PptxHandler
from src.attachments.extract.handlers.pdf import PdfHandler
from src.attachments.extract.handlers.image import ImageHandler
from src.attachments.extract.ocr.base import OcrResult
from src.attachments.extract.result import Status


class TestPlaintext(unittest.TestCase):
    def test_can_handle_and_extract(self):
        h = PlaintextHandler()
        self.assertTrue(h.can_handle("text/plain", "a.txt"))
        self.assertTrue(h.can_handle("text/csv", "a.csv"))
        self.assertTrue(h.can_handle("text/calendar", "a.ics"))
        self.assertFalse(h.can_handle("image/png", "a.png"))
        r = h.extract(b"hello world", "text/plain", "a.txt")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("hello world", r.text)

    def test_csv_extracts(self):
        r = PlaintextHandler().extract(b"a,b\n1,2", "text/csv", "x.csv")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("a,b", r.text)


class TestHtml(unittest.TestCase):
    def test_strips_tags(self):
        h = HtmlHandler()
        self.assertTrue(h.can_handle("text/html", "a.html"))
        r = h.extract(b"<p>Hi <b>there</b></p>", "text/html", "a.html")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("Hi", r.text)
        self.assertNotIn("<p>", r.text)

    def test_can_handle_negatives_and_htm(self):
        h = HtmlHandler()
        self.assertTrue(h.can_handle("", "page.htm"))
        self.assertFalse(h.can_handle("text/plain", "a.txt"))
        self.assertFalse(h.can_handle("image/png", "a.png"))


class TestOfficeHandlers(unittest.TestCase):
    def test_docx(self):
        import docx
        d = docx.Document(); d.add_paragraph("Quarterly report up 12%")
        b = io.BytesIO(); d.save(b)
        h = DocxHandler()
        self.assertTrue(h.can_handle(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "r.docx"))
        r = h.extract(b.getvalue(), "application/octet-stream", "r.docx")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("Quarterly report", r.text)

    def test_xlsx(self):
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active; ws.append(["Item", "Cost"]); ws.append(["Widget", 42])
        b = io.BytesIO(); wb.save(b)
        r = XlsxHandler().extract(b.getvalue(), "application/octet-stream", "s.xlsx")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("Widget", r.text)

    def test_pptx(self):
        import pptx
        p = pptx.Presentation(); s = p.slides.add_slide(p.slide_layouts[5])
        s.shapes.title.text = "Roadmap 2026"
        b = io.BytesIO(); p.save(b)
        r = PptxHandler().extract(b.getvalue(), "application/octet-stream", "d.pptx")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("Roadmap", r.text)

    def test_docx_missing_lib_is_binary(self):
        with mock.patch.dict("sys.modules", {"docx": None}):
            r = DocxHandler().extract(b"PK\x03\x04", "application/octet-stream", "r.docx")
        self.assertEqual(r.status, Status.BINARY)

    def test_xlsx_missing_lib_is_binary(self):
        with mock.patch.dict("sys.modules", {"openpyxl": None}):
            r = XlsxHandler().extract(b"PK\x03\x04", "application/octet-stream", "s.xlsx")
        self.assertEqual(r.status, Status.BINARY)

    def test_pptx_missing_lib_is_binary(self):
        with mock.patch.dict("sys.modules", {"pptx": None}):
            r = PptxHandler().extract(b"PK\x03\x04", "application/octet-stream", "d.pptx")
        self.assertEqual(r.status, Status.BINARY)

    def test_can_handle_xlsx_pptx_by_mime(self):
        self.assertTrue(XlsxHandler().can_handle(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "s.xlsx"))
        self.assertTrue(PptxHandler().can_handle(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation", "d.pptx"))
        self.assertFalse(XlsxHandler().can_handle("text/plain", "a.txt"))


class _StubOcr:
    def __init__(self, result):
        self.result = result
        self.calls = 0

    def read(self, data, mime, filename):
        self.calls += 1
        return self.result


class TestPdfImageHandlers(unittest.TestCase):
    def test_image_delegates_to_ocr(self):
        ocr = _StubOcr(OcrResult("scanned text", Status.EXTRACTED, "stub"))
        h = ImageHandler(ocr)
        self.assertTrue(h.can_handle("image/png", "x.png"))
        r = h.extract(b"\x89PNG", "image/png", "x.png")
        self.assertEqual(r.text, "scanned text")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertEqual(r.extractor, "stub")
        self.assertEqual(ocr.calls, 1)

    def test_pdf_textlayer_short_circuits_ocr(self):
        ocr = _StubOcr(OcrResult("", Status.OCR_UNAVAILABLE, "stub"))
        with mock.patch("src.attachments.extract.handlers.pdf._pdf_text", return_value="real text"):
            r = PdfHandler(ocr).extract(b"%PDF", "application/pdf", "x.pdf")
        self.assertEqual(r.text, "real text")
        self.assertEqual(ocr.calls, 0)

    def test_pdf_no_textlayer_delegates_to_ocr(self):
        ocr = _StubOcr(OcrResult("ocr text", Status.EXTRACTED, "stub"))
        with mock.patch("src.attachments.extract.handlers.pdf._pdf_text", return_value=""):
            r = PdfHandler(ocr).extract(b"%PDF", "application/pdf", "x.pdf")
        self.assertEqual(r.text, "ocr text")
        self.assertEqual(ocr.calls, 1)

    def test_pdf_unavailable_ocr_maps_to_ocr_unavailable(self):
        ocr = _StubOcr(OcrResult("", Status.OCR_UNAVAILABLE, "stub"))
        with mock.patch("src.attachments.extract.handlers.pdf._pdf_text", return_value=""):
            r = PdfHandler(ocr).extract(b"%PDF", "application/pdf", "x.pdf")
        self.assertEqual(r.status, Status.OCR_UNAVAILABLE)
        self.assertIn("pdf+", r.extractor)

    def test_image_can_handle_negative(self):
        h = ImageHandler(_StubOcr(OcrResult("", Status.EMPTY, "stub")))
        self.assertFalse(h.can_handle("application/pdf", "x.pdf"))
        self.assertFalse(h.can_handle("text/plain", "a.txt"))


class TestCharsetHandling(unittest.TestCase):
    """Declared charsets travel in the mime string (text/plain; charset=...) and the
    text handlers must honour them instead of mojibake-ing through latin-1."""

    SJIS = "こんにちは".encode("shift_jis")

    def test_plaintext_uses_declared_charset(self):
        r = PlaintextHandler().extract(self.SJIS, "text/plain; charset=shift_jis", "a.txt")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("こんにちは", r.text)

    def test_plaintext_bad_declared_charset_falls_back(self):
        r = PlaintextHandler().extract(b"hello", "text/plain; charset=bogus-charset", "a.txt")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("hello", r.text)

    def test_html_uses_declared_charset(self):
        html = "<p>こんにちは</p>".encode("shift_jis")
        r = HtmlHandler().extract(html, "text/html; charset=shift_jis", "a.html")
        self.assertEqual(r.status, Status.EXTRACTED)
        self.assertIn("こんにちは", r.text)

    def test_can_handle_tolerates_mime_parameters(self):
        self.assertTrue(PlaintextHandler().can_handle("text/plain; charset=utf-8", "noext"))
        self.assertTrue(HtmlHandler().can_handle("text/html; charset=iso-8859-1", "noext"))

    def test_decode_text_is_public(self):
        from src.attachments.extract.handlers.plaintext import decode_text
        self.assertEqual(decode_text("café".encode("utf-8")), "café")
        self.assertEqual(decode_text("é".encode("latin-1")), "é")
        self.assertEqual(decode_text("é".encode("latin-1"), charset="iso-8859-1"), "é")


if __name__ == "__main__":
    unittest.main()
