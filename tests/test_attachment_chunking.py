"""Tests for structure-aware attachment chunking (issue #89).

Covers ``src/indexing/attachment_chunking.py`` (the per-MIME chunker) and its
integration into ``build_attachment_documents``: spreadsheets chunk by row-groups
with the header repeated and no chunk over the token budget, PDFs by page, PPTX by
slide, and small attachments stay a single chunk (no #80 regression)."""

import importlib.util
import io
import os
import tempfile
import unittest
from email.message import EmailMessage

from src.indexing.attachment_chunking import (
    DEFAULT_CHUNK_BUDGET,
    _default_token_len,
    chunk_attachment,
)


def _have(module: str) -> bool:
    try:
        return importlib.util.find_spec(module) is not None
    except (ImportError, ValueError):
        return False


_HAVE_OPENPYXL = _have("openpyxl")
_HAVE_PANDAS = _have("pandas")
_HAVE_PPTX = _have("pptx")
_HAVE_DOCX = _have("docx")


def _xlsx_bytes(rows) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _docx_bytes(paragraphs) -> bytes:
    """A .docx from ``(style, text)`` pairs. ``style`` may be a heading style name
    (e.g. "Heading 1") or None for a normal paragraph."""
    import docx

    d = docx.Document()
    for style, text in paragraphs:
        if style:
            d.add_paragraph(text, style=style)
        else:
            d.add_paragraph(text)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _pptx_bytes(slide_texts) -> bytes:
    import pptx

    pres = pptx.Presentation()
    blank = pres.slide_layouts[6]  # blank layout
    for text in slide_texts:
        slide = pres.slides.add_slide(blank)
        box = slide.shapes.add_textbox(0, 0, 100, 100)
        box.text_frame.text = text
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


class _StubReader:
    """A pypdf-like reader whose pages return fixed text — lets the PDF path be
    exercised without a real PDF binary or an OCR backend."""

    class _Page:
        def __init__(self, text):
            self._t = text

        def extract_text(self):
            return self._t

    def __init__(self, texts):
        self.pages = [self._Page(t) for t in texts]


XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


class TestTokenLen(unittest.TestCase):
    def test_counts_numeric_and_tabular_units(self):
        # "210,000,000" fragments into digit runs + commas, mirroring a sub-word
        # tokenizer, so tabular/numeric text is never under-counted.
        self.assertGreater(_default_token_len("210,000,000"), 1)
        self.assertEqual(_default_token_len(""), 0)


@unittest.skipUnless(_HAVE_PANDAS and _HAVE_OPENPYXL, "pandas/openpyxl required")
class TestSpreadsheetChunking(unittest.TestCase):
    def test_many_rows_split_into_budgeted_chunks_with_repeated_header(self):
        rows = [["id", "name", "amount"]]
        for i in range(500):
            rows.append([i, f"customer_{i}", 1000 + i])
        data = _xlsx_bytes(rows)
        budget = 64

        chunks = chunk_attachment(
            data=data, text="", mime=XLSX_MIME, filename="big.xlsx", budget=budget
        )

        # Multiple chunks (the sheet exceeds one budget).
        self.assertGreater(len(chunks), 1)
        # Every chunk is under the token budget...
        for c in chunks:
            self.assertLessEqual(_default_token_len(c), budget)
        # ...carries the header row (self-describing mini-table)...
        for c in chunks:
            self.assertIn("id\tname\tamount", c)
        # ...and no data row is split across chunks: a known cell is in exactly one.
        hits = [c for c in chunks if "customer_237\t1237" in c]
        self.assertEqual(len(hits), 1, "row must live intact in exactly one chunk")
        # The tail row survives (the #87 truncation would have dropped it).
        self.assertTrue(any("customer_499\t1499" in c for c in chunks))

    def test_small_sheet_stays_single_chunk(self):
        # No-regression: the #80 fixture (tiny sheet) must remain one chunk carrying
        # the known value.
        data = _xlsx_bytes(
            [
                ["Objective", "Weight", "Team Target FY2025"],
                ["Partner bookings", 0.50, "210,000,000"],
            ]
        )
        chunks = chunk_attachment(
            data=data, text="", mime=XLSX_MIME, filename="Q3.xlsx", budget=DEFAULT_CHUNK_BUDGET
        )
        self.assertEqual(len(chunks), 1)
        self.assertIn("210,000,000", chunks[0])

    def test_multiple_sheets_each_get_their_own_header(self):
        import openpyxl

        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Alpha"
        ws1.append(["a", "b"])
        ws1.append([1, 2])
        ws2 = wb.create_sheet("Beta")
        ws2.append(["x", "y"])
        ws2.append([9, 8])
        buf = io.BytesIO()
        wb.save(buf)

        chunks = chunk_attachment(data=buf.getvalue(), text="", mime=XLSX_MIME, filename="two.xlsx")
        joined = "\n---\n".join(chunks)
        self.assertIn("[sheet: Alpha]", joined)
        self.assertIn("[sheet: Beta]", joined)
        self.assertTrue(any("a\tb" in c and "1\t2" in c for c in chunks))
        self.assertTrue(any("x\ty" in c and "9\t8" in c for c in chunks))

    def test_csv_with_commas_in_values_keeps_the_value(self):
        # The #80 canary: "210,000,000" survives the CSV parse/rejoin.
        data = b"item,amount\nDeal,210,000,000\n"
        chunks = chunk_attachment(
            data=data, text="item\tamount\nDeal\t210,000,000", mime="text/csv", filename="l.csv"
        )
        self.assertTrue(any("210,000,000" in c for c in chunks))

    def test_row_wider_than_budget_is_kept_intact(self):
        # A single row wider than the budget is never cut mid-row — it is emitted on
        # its own (with the header) rather than split, so the record stays coherent.
        wide = " ".join(f"c{i}" for i in range(60))
        data = _xlsx_bytes([["h"], [wide]])
        chunks = chunk_attachment(
            data=data, text="", mime=XLSX_MIME, filename="wide.xlsx", budget=10
        )
        self.assertTrue(any(wide.replace(" ", "") in c.replace(" ", "") for c in chunks))

    def test_header_only_sheet_yields_one_chunk(self):
        data = _xlsx_bytes([["only", "header", "row"]])
        chunks = chunk_attachment(data=data, text="", mime=XLSX_MIME, filename="hdr.xlsx")
        self.assertEqual(len(chunks), 1)
        self.assertIn("only\theader\trow", chunks[0])

    def test_injectable_token_len_is_used(self):
        # A custom token_len (char count) produces a different split than the default
        # word/punctuation counter — proving the callable is actually consulted.
        rows = [["k", "v"]] + [[i, i] for i in range(50)]
        data = _xlsx_bytes(rows)
        by_chars = chunk_attachment(
            data=data, text="", mime=XLSX_MIME, filename="t.xlsx", budget=120, token_len=len
        )
        by_default = chunk_attachment(
            data=data, text="", mime=XLSX_MIME, filename="t.xlsx", budget=120
        )
        # Char-counting hits the budget sooner than word-counting -> more chunks.
        # (The exact char length can exceed the budget by the inter-row newlines,
        # which the token_len sums separately; the point is the callable is used.)
        self.assertGreater(len(by_chars), len(by_default))


class TestPdfChunking(unittest.TestCase):
    def test_one_chunk_per_page(self):
        # Patch pypdf.PdfReader so the page-splitting logic is tested without a real
        # PDF/OCR. Three pages -> three chunks in reading order.
        from unittest import mock

        import pypdf

        pages = ["Page one text.", "Page two text.", "Page three text."]
        with mock.patch.object(pypdf, "PdfReader", lambda _b: _StubReader(pages)):
            chunks = chunk_attachment(
                data=b"%PDF-1.4 fake", text="ignored", mime="application/pdf", filename="doc.pdf"
            )
        self.assertEqual(len(chunks), 3)
        self.assertEqual(chunks[0], "Page one text.")
        self.assertEqual(chunks[2], "Page three text.")

    def test_over_budget_page_is_prose_split(self):
        from unittest import mock

        import pypdf

        big_page = " ".join(f"word{i}" for i in range(200))
        with mock.patch.object(pypdf, "PdfReader", lambda _b: _StubReader([big_page])):
            chunks = chunk_attachment(
                data=b"%PDF fake",
                text="ignored",
                mime="application/pdf",
                filename="d.pdf",
                budget=32,
            )
        self.assertGreater(len(chunks), 1)
        for c in chunks:
            self.assertLessEqual(_default_token_len(c), 32)


@unittest.skipUnless(_HAVE_PPTX, "python-pptx required")
class TestPptxChunking(unittest.TestCase):
    def test_one_chunk_per_slide(self):
        data = _pptx_bytes(["Slide A content", "Slide B content", "Slide C content"])
        chunks = chunk_attachment(
            data=data,
            text="ignored",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="deck.pptx",
        )
        self.assertEqual(len(chunks), 3)
        self.assertIn("Slide A content", chunks[0])
        self.assertIn("Slide C content", chunks[2])


DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


@unittest.skipUnless(_HAVE_DOCX, "python-docx required")
class TestDocxChunking(unittest.TestCase):
    def test_split_by_heading_section(self):
        data = _docx_bytes(
            [
                ("Heading 1", "Introduction"),
                (None, "Intro body paragraph."),
                ("Heading 1", "Results"),
                (None, "Results body paragraph."),
                ("Heading 1", "Conclusion"),
                (None, "Conclusion body paragraph."),
            ]
        )
        chunks = chunk_attachment(data=data, text="ignored", mime=DOCX_MIME, filename="report.docx")
        # One chunk per heading-delimited section.
        self.assertEqual(len(chunks), 3)
        self.assertIn("Introduction", chunks[0])
        self.assertIn("Intro body paragraph.", chunks[0])
        self.assertIn("Results", chunks[1])
        self.assertIn("Conclusion", chunks[2])

    def test_over_budget_section_is_prose_split(self):
        long_body = " ".join(f"word{i}" for i in range(200))
        data = _docx_bytes([("Heading 1", "Big section"), (None, long_body)])
        chunks = chunk_attachment(
            data=data, text="ignored", mime=DOCX_MIME, filename="big.docx", budget=32
        )
        self.assertGreater(len(chunks), 1)
        for c in chunks:
            self.assertLessEqual(_default_token_len(c), 32)

    def test_headingless_docx_is_single_prose_split(self):
        data = _docx_bytes([(None, "Just a plain paragraph with no headings.")])
        chunks = chunk_attachment(data=data, text="ignored", mime=DOCX_MIME, filename="plain.docx")
        self.assertEqual(len(chunks), 1)
        self.assertIn("plain paragraph", chunks[0])


class TestFallback(unittest.TestCase):
    def test_prose_text_falls_back_to_budgeted_prose_chunks(self):
        text = "\n".join(f"Sentence number {i} here." for i in range(100))
        chunks = chunk_attachment(
            data=b"", text=text, mime="text/plain", filename="notes.txt", budget=20
        )
        self.assertGreater(len(chunks), 1)
        for c in chunks:
            self.assertLessEqual(_default_token_len(c), 20)

    def test_unparseable_spreadsheet_bytes_fall_back_to_text(self):
        # Garbage bytes with an .xlsx name: the structured parse returns None and the
        # flattened text is used instead (no crash, no data loss).
        chunks = chunk_attachment(
            data=b"\x00\x01not-a-zip",
            text="row1\trow2\nval",
            mime=XLSX_MIME,
            filename="broken.xlsx",
            budget=DEFAULT_CHUNK_BUDGET,
        )
        self.assertTrue(chunks)
        self.assertIn("val", "\n".join(chunks))

    def test_empty_text_yields_no_chunks(self):
        self.assertEqual(chunk_attachment(data=b"", text="   ", mime="", filename="x"), [])


def _write_eml(path, message_id, body, attachments):
    m = EmailMessage()
    m["From"] = "a@x.com"
    m["To"] = "b@x.com"
    m["Subject"] = "with attachment"
    m["Message-ID"] = message_id
    m.set_content(body)
    for data, maintype, subtype, filename in attachments:
        m.add_attachment(data, maintype=maintype, subtype=subtype, filename=filename)
    with open(path, "wb") as fh:
        fh.write(bytes(m))
    return path


@unittest.skipUnless(_HAVE_PANDAS and _HAVE_OPENPYXL, "pandas/openpyxl required")
class TestBuildAttachmentDocumentsChunking(unittest.TestCase):
    """End-to-end through build_attachment_documents: a big sheet becomes MANY
    Documents, each budgeted and header-carrying, with lineage preserved."""

    def setUp(self):
        self.d = tempfile.mkdtemp()
        import shutil

        self.addCleanup(shutil.rmtree, self.d, ignore_errors=True)

    def test_large_xlsx_becomes_many_budgeted_documents_with_lineage(self):
        from src.indexing.attachment_docs import build_attachment_documents

        rows = [["id", "name", "amount"]]
        for i in range(400):
            rows.append([i, f"customer_{i}", 1000 + i])
        path = _write_eml(
            os.path.join(self.d, "big.eml"),
            "<big@x>",
            "see attached",
            [
                (
                    _xlsx_bytes(rows),
                    "application",
                    "vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "big.xlsx",
                )
            ],
        )

        budget = 64
        docs = build_attachment_documents([path], extractor_name="tesseract", chunk_budget=budget)

        self.assertGreater(len(docs), 1, "a 400-row sheet must yield multiple chunks")
        for doc in docs:
            # Budget respected, header repeated, lineage preserved on every chunk.
            self.assertLessEqual(_default_token_len(doc.text), budget)
            self.assertIn("id\tname\tamount", doc.text)
            self.assertEqual(doc.metadata["content_kind"], "attachment")
            self.assertEqual(doc.metadata["attachment_name"], "big.xlsx")
            self.assertEqual(doc.metadata["parent_message_id"], "<big@x>")
            self.assertTrue(doc.metadata.get("thread_id"))
            self.assertIn("chunk_index", doc.metadata)
        # A known cell lands in exactly one chunk (no row split).
        hits = [d for d in docs if "customer_237\t1237" in d.text]
        self.assertEqual(len(hits), 1)

    def test_small_xlsx_stays_single_document(self):
        # #80 no-regression at the Document level: one small sheet -> one Document.
        from src.indexing.attachment_docs import build_attachment_documents

        path = _write_eml(
            os.path.join(self.d, "small.eml"),
            "<mb01@northwind.example>",
            "Team\nHere are MBO targets\nLMK",
            [
                (
                    _xlsx_bytes(
                        [
                            ["Objective", "Weight", "Team Target FY2025"],
                            ["Partner bookings", 0.50, "210,000,000"],
                        ]
                    ),
                    "application",
                    "vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "Q3 MBO targets partner team.xlsx",
                )
            ],
        )
        docs = build_attachment_documents([path], extractor_name="tesseract")
        self.assertEqual(len(docs), 1)
        self.assertIn("210,000,000", docs[0].text)
        # Single-chunk attachments are not chunk-index tagged.
        self.assertNotIn("chunk_index", docs[0].metadata)


if __name__ == "__main__":
    unittest.main()
